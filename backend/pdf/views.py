import os
import io
import base64
from datetime import datetime
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_http_methods
from django.core.files.uploadedfile import InMemoryUploadedFile
from django.conf import settings
import tempfile

# PDF processing imports
import PyPDF2
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

# Try to import advanced libraries
try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False


@csrf_exempt
@require_http_methods(["POST"])
def merge_pdfs(request):
    """Merge multiple PDF files into one"""
    try:
        files = request.FILES.getlist('pdfs')
        if len(files) < 2:
            return JsonResponse({'error': 'At least 2 PDF files are required'}, status=400)
        
        merger = PyPDF2.PdfMerger()
        
        for pdf_file in files:
            merger.append(pdf_file)
        
        # Save merged PDF to buffer
        pdf_buffer = io.BytesIO()
        merger.write(pdf_buffer)
        merger.close()
        pdf_buffer.seek(0)
        
        # Encode to base64
        pdf_base64 = base64.b64encode(pdf_buffer.getvalue()).decode()
        
        return JsonResponse({
            'success': True,
            'pdf': f'data:application/pdf;base64,{pdf_base64}',
            'filename': 'merged.pdf'
        })
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
@require_http_methods(["POST"])
def split_pdf(request):
    """Split PDF into multiple files"""
    try:
        # Accept both 'pdf' and 'file' as valid file keys
        if 'pdf' in request.FILES:
            pdf_file = request.FILES['pdf']
        elif 'file' in request.FILES:
            pdf_file = request.FILES['file']
        else:
            return JsonResponse({'error': 'No PDF file provided'}, status=400)
        split_type = request.POST.get('split_type', 'single')  # 'single' or 'range'
        page_ranges = request.POST.get('page_ranges', '')  # For range split
        
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        total_pages = len(pdf_reader.pages)
        
        result_files = []
        
        if split_type == 'single':
            # Split each page into separate PDF
            for page_num in range(total_pages):
                pdf_writer = PyPDF2.PdfWriter()
                pdf_writer.add_page(pdf_reader.pages[page_num])
                
                pdf_buffer = io.BytesIO()
                pdf_writer.write(pdf_buffer)
                pdf_buffer.seek(0)
                
                pdf_base64 = base64.b64encode(pdf_buffer.getvalue()).decode()
                result_files.append({
                    'filename': f'page_{page_num + 1}.pdf',
                    'pdf': f'data:application/pdf;base64,{pdf_base64}'
                })
        
        elif split_type == 'range' and page_ranges:
            # Split by page ranges (e.g., "1-3,5,7-9")
            ranges = [r.strip() for r in page_ranges.split(',')]
            
            for range_item in ranges:
                if '-' in range_item:
                    start, end = map(int, range_item.split('-'))
                    page_indices = list(range(start - 1, min(end, total_pages)))
                else:
                    page_num = int(range_item)
                    page_indices = [page_num - 1] if page_num <= total_pages else []
                
                if page_indices:
                    pdf_writer = PyPDF2.PdfWriter()
                    for page_idx in page_indices:
                        pdf_writer.add_page(pdf_reader.pages[page_idx])
                    
                    pdf_buffer = io.BytesIO()
                    pdf_writer.write(pdf_buffer)
                    pdf_buffer.seek(0)
                    
                    pdf_base64 = base64.b64encode(pdf_buffer.getvalue()).decode()
                    result_files.append({
                        'filename': f'pages_{range_item.replace("-", "_")}.pdf',
                        'pdf': f'data:application/pdf;base64,{pdf_base64}'
                    })
        
        return JsonResponse({
            'success': True,
            'files': result_files,
            'total_pages': total_pages
        })
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
@require_http_methods(["POST"])
def pdf_to_images(request):
    """Convert PDF pages to images"""
    try:
        if not PYMUPDF_AVAILABLE:
            return JsonResponse({
                'error': 'PDF to images conversion requires PyMuPDF. Install with: pip install PyMuPDF'
            }, status=501)
            
        # Accept both 'pdf' and 'file' as valid file keys
        if 'pdf' in request.FILES:
            pdf_file = request.FILES['pdf']
        elif 'file' in request.FILES:
            pdf_file = request.FILES['file']
        else:
            return JsonResponse({'error': 'No PDF file provided'}, status=400)
        image_format = request.POST.get('format', 'PNG').upper()
        dpi = int(request.POST.get('dpi', 150))
        
        if image_format not in ['PNG', 'JPEG', 'WEBP']:
            return JsonResponse({'error': 'Unsupported format. Use PNG, JPEG, or WEBP'}, status=400)
        
        # Open PDF with PyMuPDF
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        
        result_images = []
        
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            
            # Render page to image
            mat = fitz.Matrix(dpi / 72, dpi / 72)
            pix = page.get_pixmap(matrix=mat)
            
            # Convert to base64
            img_data = pix.tobytes(image_format.lower())
            img_base64 = base64.b64encode(img_data).decode()
            
            result_images.append({
                'page': page_num + 1,
                'image': f'data:image/{image_format.lower()};base64,{img_base64}',
                'filename': f'page_{page_num + 1}.{image_format.lower()}'
            })
        
        pdf_document.close()
        
        return JsonResponse({
            'success': True,
            'images': result_images,
            'total_pages': len(result_images)
        })
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
@require_http_methods(["POST"])
def pdf_to_word(request):
    """Convert PDF to Word document"""
    try:
        if not DOCX_AVAILABLE:
            return JsonResponse({
                'error': 'PDF to Word conversion requires python-docx. Install with: pip install python-docx'
            }, status=501)
            
        if not PYMUPDF_AVAILABLE:
            return JsonResponse({
                'error': 'PDF to Word conversion requires PyMuPDF. Install with: pip install PyMuPDF'
            }, status=501)
            
        # Accept both 'pdf' and 'file' as valid file keys
        if 'pdf' in request.FILES:
            pdf_file = request.FILES['pdf']
        elif 'file' in request.FILES:
            pdf_file = request.FILES['file']
        else:
            return JsonResponse({'error': 'No PDF file provided'}, status=400)
        
        # Open PDF with PyMuPDF
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        
        # Create Word document
        doc = Document()
        
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            text = page.get_text()
            
            if text.strip():
                doc.add_heading(f'Page {page_num + 1}', level=2)
                doc.add_paragraph(text)
                doc.add_page_break()
        
        pdf_document.close()
        
        # Save Word document to buffer
        doc_buffer = io.BytesIO()
        doc.save(doc_buffer)
        doc_buffer.seek(0)
        
        # Encode to base64
        doc_base64 = base64.b64encode(doc_buffer.getvalue()).decode()
        
        return JsonResponse({
            'success': True,
            'docx': f'data:application/vnd.openxmlformats-officedocument.wordprocessingml.document;base64,{doc_base64}',
            'filename': 'converted.docx'
        })
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
@require_http_methods(["POST"])
def pdf_to_powerpoint(request):
    """Convert PDF to PowerPoint"""
    try:
        # Accept both 'pdf' and 'file' as valid file keys
        if 'pdf' in request.FILES:
            pdf_file = request.FILES['pdf']
        elif 'file' in request.FILES:
            pdf_file = request.FILES['file']
        else:
            return JsonResponse({'error': 'No PDF file provided'}, status=400)
        
        # Open PDF with PyMuPDF
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            
            # Convert page to image
            mat = fitz.Matrix(2, 2)  # Higher resolution for slides
            pix = page.get_pixmap(matrix=mat)
            
            # Add slide
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add image to slide
            img_data = pix.tobytes("png")
            img_buffer = io.BytesIO(img_data)
            slide.shapes.add_picture(img_buffer, Inches(0.5), Inches(0.5), width=Inches(9))
        
        pdf_document.close()
        
        # Save PowerPoint to buffer
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        # Encode to base64
        ppt_base64 = base64.b64encode(ppt_buffer.getvalue()).decode()
        
        return JsonResponse({
            'success': True,
            'pptx': f'data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{ppt_base64}',
            'filename': 'converted.pptx'
        })
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
@require_http_methods(["POST"])
def pdf_to_excel(request):
    """Extract tables from PDF to Excel"""
    try:
        # Accept both 'pdf' and 'file' as valid file keys
        if 'pdf' in request.FILES:
            pdf_file = request.FILES['pdf']
        elif 'file' in request.FILES:
            pdf_file = request.FILES['file']
        else:
            return JsonResponse({'error': 'No PDF file provided'}, status=400)
        
        # Open PDF with PyMuPDF
        pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
        
        # Create Excel workbook
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Extracted Data"
        
        row_num = 1
        
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            text = page.get_text()
            
            if text.strip():
                # Add page header
                ws.cell(row=row_num, column=1, value=f"Page {page_num + 1}")
                row_num += 1
                
                # Split text into lines and add to Excel
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        ws.cell(row=row_num, column=1, value=line.strip())
                        row_num += 1
                
                # Add empty row between pages
                ws.cell(row=row_num, column=1, value="")
                row_num += 1
        
        pdf_document.close()
        
        # Save Excel to buffer
        excel_buffer = io.BytesIO()
        wb.save(excel_buffer)
        excel_buffer.seek(0)
        
        # Encode to base64
        excel_base64 = base64.b64encode(excel_buffer.getvalue()).decode()
        
        return JsonResponse({
            'success': True,
            'xlsx': f'data:application/vnd.openxmlformats-officedocument.spreadsheetml.sheet;base64,{excel_base64}',
            'filename': 'extracted.xlsx'
        })
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
@require_http_methods(["POST"])
def protect_pdf(request):
    """Add password protection to PDF"""
    try:
        # Accept both 'pdf' and 'file' as valid file keys
        if 'pdf' in request.FILES:
            pdf_file = request.FILES['pdf']
        elif 'file' in request.FILES:
            pdf_file = request.FILES['file']
        else:
            return JsonResponse({'error': 'No PDF file provided'}, status=400)
        password = request.POST.get('password', '')
        
        if not password:
            return JsonResponse({'error': 'Password is required'}, status=400)
        
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        pdf_writer = PyPDF2.PdfWriter()
        
        # Add all pages to writer
        for page in pdf_reader.pages:
            pdf_writer.add_page(page)
        
        # Add password protection
        pdf_writer.encrypt(password)
        
        # Save protected PDF to buffer
        pdf_buffer = io.BytesIO()
        pdf_writer.write(pdf_buffer)
        pdf_buffer.seek(0)
        
        # Encode to base64
        pdf_base64 = base64.b64encode(pdf_buffer.getvalue()).decode()
        
        return JsonResponse({
            'success': True,
            'pdf': f'data:application/pdf;base64,{pdf_base64}',
            'filename': 'protected.pdf'
        })
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
@require_http_methods(["POST"])
def unlock_pdf(request):
    """Remove password protection from PDF"""
    try:
        # Accept both 'pdf' and 'file' as valid file keys
        if 'pdf' in request.FILES:
            pdf_file = request.FILES['pdf']
        elif 'file' in request.FILES:
            pdf_file = request.FILES['file']
        else:
            return JsonResponse({'error': 'No PDF file provided'}, status=400)
        password = request.POST.get('password', '')
        
        if not password:
            return JsonResponse({'error': 'Password is required to unlock PDF'}, status=400)
        
        try:
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            if pdf_reader.is_encrypted:
                pdf_reader.decrypt(password)
        except:
            return JsonResponse({'error': 'Incorrect password or corrupted PDF'}, status=400)
        
        pdf_writer = PyPDF2.PdfWriter()
        
        # Add all pages to writer
        for page in pdf_reader.pages:
            pdf_writer.add_page(page)
        
        # Save unlocked PDF to buffer
        pdf_buffer = io.BytesIO()
        pdf_writer.write(pdf_buffer)
        pdf_buffer.seek(0)
        
        # Encode to base64
        pdf_base64 = base64.b64encode(pdf_buffer.getvalue()).decode()
        
        return JsonResponse({
            'success': True,
            'pdf': f'data:application/pdf;base64,{pdf_base64}',
            'filename': 'unlocked.pdf'
        })
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
@require_http_methods(["POST"])
def rotate_pdf(request):
    """Rotate PDF pages"""
    try:
        # Accept both 'pdf' and 'file' as valid file keys
        if 'pdf' in request.FILES:
            pdf_file = request.FILES['pdf']
        elif 'file' in request.FILES:
            pdf_file = request.FILES['file']
        else:
            return JsonResponse({'error': 'No PDF file provided'}, status=400)
        rotation = int(request.POST.get('rotation', 90))  # 90, 180, or 270 degrees
        
        if rotation not in [90, 180, 270]:
            return JsonResponse({'error': 'Rotation must be 90, 180, or 270 degrees'}, status=400)
        
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        pdf_writer = PyPDF2.PdfWriter()
        
        # Rotate all pages
        for page in pdf_reader.pages:
            page.rotate(rotation)
            pdf_writer.add_page(page)
        
        # Save rotated PDF to buffer
        pdf_buffer = io.BytesIO()
        pdf_writer.write(pdf_buffer)
        pdf_buffer.seek(0)
        
        # Encode to base64
        pdf_base64 = base64.b64encode(pdf_buffer.getvalue()).decode()
        
        return JsonResponse({
            'success': True,
            'pdf': f'data:application/pdf;base64,{pdf_base64}',
            'filename': f'rotated_{rotation}.pdf'
        })
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
@require_http_methods(["POST"])
def remove_pages(request):
    """Remove specific pages from PDF"""
    try:
        # Accept both 'pdf' and 'file' as valid file keys
        if 'pdf' in request.FILES:
            pdf_file = request.FILES['pdf']
        elif 'file' in request.FILES:
            pdf_file = request.FILES['file']
        else:
            return JsonResponse({'error': 'No PDF file provided'}, status=400)
        pages_to_remove = request.POST.get('pages_to_remove', '')
        
        if not pages_to_remove:
            return JsonResponse({'error': 'Pages to remove is required'}, status=400)
        
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        total_pages = len(pdf_reader.pages)
        
        # Parse pages to remove (e.g., "1,3,5-7")
        pages_set = set()
        for item in pages_to_remove.split(','):
            item = item.strip()
            if '-' in item:
                start, end = map(int, item.split('-'))
                pages_set.update(range(start, min(end + 1, total_pages + 1)))
            else:
                page_num = int(item)
                if 1 <= page_num <= total_pages:
                    pages_set.add(page_num)
        
        pdf_writer = PyPDF2.PdfWriter()
        
        # Add all pages except those to remove
        for page_num in range(1, total_pages + 1):
            if page_num not in pages_set:
                pdf_writer.add_page(pdf_reader.pages[page_num - 1])
        
        # Save PDF to buffer
        pdf_buffer = io.BytesIO()
        pdf_writer.write(pdf_buffer)
        pdf_buffer.seek(0)
        
        # Encode to base64
        pdf_base64 = base64.b64encode(pdf_buffer.getvalue()).decode()
        
        return JsonResponse({
            'success': True,
            'pdf': f'data:application/pdf;base64,{pdf_base64}',
            'filename': 'pages_removed.pdf',
            'original_pages': total_pages,
            'new_pages': total_pages - len(pages_set)
        })
        
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)
